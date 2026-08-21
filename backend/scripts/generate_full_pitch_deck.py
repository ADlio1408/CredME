import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_filename="CredMe-Presentation.pptx"):
    prs = Presentation()
    # Set 16:9 widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette Constants
    C_NAVY_DARK   = RGBColor(15, 23, 42)     # #0F172A (Slide Backgrounds / Dark headers)
    C_NAVY_CARD   = RGBColor(30, 41, 59)     # #1E293B (Dark Cards)
    C_SLATE_BG    = RGBColor(248, 250, 252)  # #F8FAFC (Light Slide Background)
    C_WHITE       = RGBColor(255, 255, 255)  # #FFFFFF (Card Backgrounds / White Text)
    C_PRIMARY     = RGBColor(79, 70, 229)    # #4F46E5 (Indigo Accent)
    C_PRIMARY_BG  = RGBColor(238, 242, 255)  # #EEF2FF (Light Indigo Card)
    C_PURPLE      = RGBColor(124, 58, 237)   # #7C3AED (Violet Accent)
    C_TEXT_DARK   = RGBColor(15, 23, 42)     # #0F172A (Primary Text Dark)
    C_TEXT_MUTED  = RGBColor(100, 116, 139)  # #64748B (Secondary Text)
    C_TEXT_LIGHT  = RGBColor(241, 245, 249)  # #F1F5F9 (Primary Text Light)
    C_GREEN       = RGBColor(16, 185, 129)   # #10B981 (Success / Positive)
    C_GREEN_BG    = RGBColor(236, 253, 245)  # #ECFDF5 (Light Green Card)
    C_RED         = RGBColor(239, 68, 68)    # #EF4444 (Alert / Negative)
    C_RED_BG      = RGBColor(254, 242, 242)  # #FEF2F2 (Light Red Card)
    C_BORDER      = RGBColor(226, 232, 240)  # #E2E8F0 (Card Border)

    def set_slide_background(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category="CREDME — NEXT-GEN CREDIT INTELLIGENCE", is_dark=False):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = C_PRIMARY if not is_dark else RGBColor(165, 180, 252)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.55))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = C_TEXT_DARK if not is_dark else C_WHITE

    def add_card(slide, left, top, width, height, bg_color=C_WHITE, border_color=C_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()
        return card

    # ==========================================
    # SLIDE 1: TITLE SLIDE (Dark Premium Theme)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, C_NAVY_DARK)

    # Accent decorative banner card
    dec_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    dec_card.fill.solid()
    dec_card.fill.fore_color.rgb = C_NAVY_CARD
    dec_card.line.color.rgb = RGBColor(51, 65, 85)
    dec_card.line.width = Pt(1.5)

    # Category Pill
    pill = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(1.3), Inches(5.8), Inches(0.4))
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(49, 46, 129)
    pill.line.color.rgb = C_PRIMARY
    tf_pill = pill.text_frame
    p_pill = tf_pill.paragraphs[0]
    p_pill.text = "SYNCHRONY HACKATHON SOLUTION PITCH"
    p_pill.font.size = Pt(11)
    p_pill.font.bold = True
    p_pill.font.color.rgb = RGBColor(199, 210, 254)
    p_pill.alignment = PP_ALIGN.CENTER

    # Main Title
    tbox = s1.shapes.add_textbox(Inches(1.3), Inches(1.9), Inches(10.7), Inches(1.8))
    tf = tbox.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "CredMe"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = C_WHITE

    p2 = tf.add_paragraph()
    p2.text = "Next-Gen Real-Time, Multi-Modal Credit Underwriting Engine"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(165, 180, 252)

    # Subtitle / Hook
    sbox = s1.shapes.add_textbox(Inches(1.3), Inches(3.8), Inches(10.7), Inches(1.2))
    stf = sbox.text_frame
    stf.word_wrap = True
    sp1 = stf.paragraphs[0]
    sp1.text = "Expanding financial inclusion for New-to-Credit (NTC) and Thin-File applicants via fused XGBoost credit scoring, real-time behavioral anomaly detection, alternative signals, and regulatory-grade explainability."
    sp1.font.size = Pt(13)
    sp1.font.color.rgb = RGBColor(203, 213, 225)

    # Footer Card Info
    foot = s1.shapes.add_textbox(Inches(1.3), Inches(5.4), Inches(10.7), Inches(0.8))
    ftf = foot.text_frame
    fp = ftf.paragraphs[0]
    fp.text = "Problem Statement: Next-Gen Credit Intelligence | Submission Format: PPT / Pitch Deck"
    fp.font.size = Pt(12)
    fp.font.bold = True
    fp.font.color.rgb = RGBColor(148, 163, 184)

    fp2 = ftf.add_paragraph()
    fp2.text = "Candidate Roll Number: [Your Roll Number] | GitHub: https://github.com/ADlio1408/CredME"
    fp2.font.size = Pt(12)
    fp2.font.color.rgb = RGBColor(129, 140, 248)


    # ==========================================
    # SLIDE 2: THE CHALLENGE & EXECUTIVE SUMMARY
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, C_SLATE_BG)
    add_header(s2, "The Challenge: Breaking the Traditional Underwriting Paradox", "1. Problem & Executive Summary")

    # Left Column: The Problem (Red accented card)
    add_card(s2, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.5), bg_color=C_WHITE, border_color=RGBColor(254, 202, 202))
    p_box = s2.shapes.add_textbox(Inches(1.1), Inches(1.6), Inches(5.0), Inches(5.1))
    ptf = p_box.text_frame
    ptf.word_wrap = True
    
    p = ptf.paragraphs[0]
    p.text = "Traditional Underwriting Bottlenecks"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_RED

    points_left = [
        ("The 'Thin-File' Penalty:", "Millions of young, gig-economy, and new-to-credit individuals have no credit score (CreditScore = 0). Legacy systems treat 0 as 'toxic credit' and auto-decline them."),
        ("Static & Reactive Signals:", "Bureau scores reflect historical payments from months ago. They fail to catch acute real-time distress or active fraud (e.g. sudden balance depletion, suspicious logins)."),
        ("The Black-Box Dilemma:", "Modern ML/AI models are often rejected by risk & compliance teams due to lack of auditability, regulatory transparency, and demographic bias.")
    ]
    for title, desc in points_left:
        p = ptf.add_paragraph()
        p.space_before = Pt(12)
        p.text = f"•  {title} "
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = C_TEXT_DARK
        p = ptf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = C_TEXT_MUTED

    # Right Column: The CredMe Solution (Indigo/Green card)
    add_card(s2, Inches(6.9), Inches(1.4), Inches(5.6), Inches(5.5), bg_color=C_WHITE, border_color=RGBColor(199, 210, 254))
    s_box = s2.shapes.add_textbox(Inches(7.2), Inches(1.6), Inches(5.0), Inches(5.1))
    stf = s_box.text_frame
    stf.word_wrap = True

    p = stf.paragraphs[0]
    p.text = "The CredMe Solution Architecture"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    points_right = [
        ("NTC-First Philosophy (Zero != Bad):", "Thin-file applicants are routed into dedicated financial validation workflows. They can never be auto-declined on thin-file status alone."),
        ("Multi-Modal Decision Fusion:", "Combines XGBoost credit scoring + IsolationForest transaction anomaly detection + Alternative data (RentPaymentConsistency)."),
        ("Explainable & Guardrailed AI:", "Top-5 SHAP feature impact drivers (+/- percentages), four-fifths fairness audit, and plain-English narrative explanations with compliance guardrails.")
    ]
    for title, desc in points_right:
        p = stf.add_paragraph()
        p.space_before = Pt(12)
        p.text = f"✓  {title} "
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = C_PRIMARY
        p = stf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = C_TEXT_MUTED


    # ==========================================
    # SLIDE 3: CORE PHILOSOPHY: THIN-FILE & NTC
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, C_SLATE_BG)
    add_header(s3, "Solving the Thin-File Dilemma: Fairness Without Compromising Risk", "2. Core Innovation")

    cards_data = [
        ("1. Intelligent Imputation", "During training, CreditScore=0 is treated as missing rather than low credit, imputed using the median of creditworthy peers so the model never penalizes missing bureau data.", C_PRIMARY_BG, C_PRIMARY),
        ("2. Dedicated Inference Routing", "Thin-file applicants bypass standard credit thresholds and enter a holistic rule screen evaluating DTI, liquidity, income, and debt ratios. Thin-file alone cannot decline.", C_GREEN_BG, C_GREEN),
        ("3. Alternative Signal Ingestion", "Accepts optional modern signals like RentPaymentConsistency (threshold >= 0.70). Demonstrates pluggability of utility, telecom, and rent cashflow data.", C_PRIMARY_BG, C_PURPLE),
    ]

    for i, (title, desc, bg, acc) in enumerate(cards_data):
        top_pos = Inches(1.4 + i * 1.8)
        add_card(s3, Inches(0.8), top_pos, Inches(11.7), Inches(1.6), bg_color=C_WHITE, border_color=C_BORDER)
        
        # Left color bar
        bar = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top_pos, Inches(0.2), Inches(1.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = acc
        bar.line.fill.background()

        tb = s3.shapes.add_textbox(Inches(1.2), top_pos + Inches(0.2), Inches(11.0), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = acc

        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_MUTED


    # ==========================================
    # SLIDE 4: SYSTEM ARCHITECTURE
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4, C_SLATE_BG)
    add_header(s4, "Full-Stack System Architecture & Technology Stack", "3. System Design")

    layers = [
        ("Frontend UI", "React 19 + Vite\n\n• Real-time applicant form\n• SHAP breakdown cards\n• Live transaction visualizer\n• Interactive scenario testing", C_PRIMARY),
        ("FastAPI Gateway", "Python API Service\n\n• Pydantic input validation\n• Role-based API Key Auth\n• WebSocket streaming feed\n• Modular routers & CORS", C_PURPLE),
        ("ML & AI Engine", "XGBoost + IsolationForest\n\n• SHAP TreeExplainer\n• Real-time behavioral scoring\n• LLM narrative scaffold\n• Regex compliance guardrails", C_GREEN),
        ("Data & DevOps", "Docker & Persistence\n\n• Docker Compose multi-service\n• Joblib model artifacts\n• CSV historical baselines\n• Prepared for pgvector", C_NAVY_DARK),
    ]

    for i, (title, content, color) in enumerate(layers):
        left_pos = Inches(0.8 + i * 3.0)
        add_card(s4, left_pos, Inches(1.4), Inches(2.8), Inches(5.5), bg_color=C_WHITE, border_color=C_BORDER)

        # Header tag on card
        top_tag = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(1.4), Inches(2.8), Inches(0.6))
        top_tag.fill.solid()
        top_tag.fill.fore_color.rgb = color
        top_tag.line.fill.background()
        tf_tag = top_tag.text_frame
        p = tf_tag.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER

        tb = s4.shapes.add_textbox(left_pos + Inches(0.15), Inches(2.2), Inches(2.5), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(11)
        p.font.color.rgb = C_TEXT_DARK


    # ==========================================
    # SLIDE 5: MULTI-MODAL DECISION FUSION
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5, C_SLATE_BG)
    add_header(s5, "Multi-Modal Decision Fusion: Credit + Behavioral + Rules", "4. Underwriting Engine")

    # 3 Column Process
    cols = [
        ("Step 1: Credit ML Scoring", "XGBoost evaluates debt-to-income, loan duration, income stability, and credit history.\n\nOutputs approval probability (P >= 0.70 for APPROVE threshold).", C_PRIMARY_BG, C_PRIMARY),
        ("Step 2: Behavioral Anomaly", "IsolationForest scores live transactions against account baselines:\n\n• Transaction-to-balance ratio\n• Rapid login spikes\n• Velocity anomalies", C_PRIMARY_BG, C_PURPLE),
        ("Step 3: Decision Fusion Matrix", "Fuses credit probability + behavioral risk + rule screen:\n\n• High credit + High anomaly -> REVIEW\n• Thin-file + Good DTI -> APPROVE\n• Thin-file + Bad rent -> REVIEW", C_GREEN_BG, C_GREEN),
    ]

    for i, (title, content, bg, color) in enumerate(cols):
        left_pos = Inches(0.8 + i * 4.0)
        add_card(s5, left_pos, Inches(1.4), Inches(3.7), Inches(3.8), bg_color=C_WHITE, border_color=C_BORDER)

        head = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(1.4), Inches(3.7), Inches(0.5))
        head.fill.solid()
        head.fill.fore_color.rgb = color
        head.line.fill.background()
        tf_h = head.text_frame
        p = tf_h.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER

        tb = s5.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.1), Inches(3.3), Inches(2.9))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(11)
        p.font.color.rgb = C_TEXT_DARK

    # Bottom Banner: Live Streaming
    add_card(s5, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.5), bg_color=C_NAVY_DARK, border_color=None)
    tb_b = s5.shapes.add_textbox(Inches(1.1), Inches(5.5), Inches(11.1), Inches(1.3))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    p = tf_b.paragraphs[0]
    p.text = "⚡ Real-Time Streaming Ingestion (POST /stream/transaction & WS /stream/live)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(165, 180, 252)
    p2 = tf_b.add_paragraph()
    p2.space_before = Pt(4)
    p2.text = "Allows underwriting underwriters and fraud engines to monitor incoming customer transactions live via WebSockets. Each transaction updates the customer's behavioral baseline dynamically in memory."
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(226, 232, 240)


    # ==========================================
    # SLIDE 6: RESPONSIBLE AI, FAIRNESS & SHAP
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6, C_SLATE_BG)
    add_header(s6, "Responsible AI: SHAP Explainability & Disparate Impact Auditing", "5. Trust & Transparency")

    # Left: SHAP Explanations
    add_card(s6, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.5), bg_color=C_WHITE, border_color=C_BORDER)
    tb_shap = s6.shapes.add_textbox(Inches(1.1), Inches(1.6), Inches(5.0), Inches(5.1))
    tf_shap = tb_shap.text_frame
    tf_shap.word_wrap = True
    p = tf_shap.paragraphs[0]
    p.text = "SHAP Decision Factor Attribution"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    shap_points = [
        ("Mathematical Transparency:", "Calculates Shapley values for every single prediction using TreeExplainer, isolating the exact marginal contribution of each feature."),
        ("Top-5 Normalized Impact Drivers:", "Presents normalized percentage contributions (e.g. Total Debt-to-Income: -42.3%, Loan Duration: +18.5%) so applicants understand why a decision was reached."),
        ("Adverse Action Transparency:", "Directly fulfills FCRA and CFPB requirements for adverse action notices without manual underwriter burden.")
    ]
    for t, d in shap_points:
        p = tf_shap.add_paragraph()
        p.space_before = Pt(10)
        p.text = f"•  {t}"
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = C_TEXT_DARK
        p = tf_shap.add_paragraph()
        p.text = d
        p.font.size = Pt(11)
        p.font.color.rgb = C_TEXT_MUTED

    # Right: Fairness Audit (Four-Fifths Rule)
    add_card(s6, Inches(6.9), Inches(1.4), Inches(5.6), Inches(5.5), bg_color=C_WHITE, border_color=C_BORDER)
    tb_fair = s6.shapes.add_textbox(Inches(7.2), Inches(1.6), Inches(5.0), Inches(5.1))
    tf_fair = tb_fair.text_frame
    tf_fair.word_wrap = True
    p = tf_fair.paragraphs[0]
    p.text = "Four-Fifths Fairness Audit (EEOC Standard)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_PURPLE

    fair_points = [
        ("Automated Disparate Impact Screening:", "Audits credit model approvals across Age groups (<25, 25-39, 40-59, 60+) and Marital Status via `backend/fairness_audit.py`."),
        ("Live Audit Endpoint (`GET /fairness/report`):", "Publishes pre-computed demographic parity metrics, disparate impact ratios, and reference distributions directly via the API."),
        ("Zero Bias Amplification:", "Analysis confirms predicted approval rates mirror actual historical baseline rates, validating that the ML model does not amplify historical biases.")
    ]
    for t, d in fair_points:
        p = tf_fair.add_paragraph()
        p.space_before = Pt(10)
        p.text = f"✓  {t}"
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = C_PURPLE
        p = tf_fair.add_paragraph()
        p.text = d
        p.font.size = Pt(11)
        p.font.color.rgb = C_TEXT_MUTED


    # ==========================================
    # SLIDE 7: AI NARRATIVE LAYER & GUARDRAILS
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7, C_SLATE_BG)
    add_header(s7, "Generative AI Narrative Layer & Regulatory Guardrails", "6. Generative AI Architecture")

    # 3 Horizontal Process Cards
    narrative_steps = [
        ("1. Strict Decoupling (No Autonomous LLM Decisions)", "The LLM never makes credit or risk decisions. It strictly acts as a natural language synthesizer, translating already-computed SHAP values and rule outputs into human-readable customer narratives (`POST /explain/narrative`).", C_PRIMARY),
        ("2. Multi-Tier Compliance Guardrails", "Automated pre/post-generation regex and policy filters screen output for prohibited protected characteristics (race, religion, gender, medical status) and guarantee no false commitment to overturn decisions.", C_PURPLE),
        ("3. Offline-First Template Fallback Architecture", "If no OpenAI API key is provided (`CREDME_LLM_API_KEY`), the engine automatically switches to a deterministic template fallback (`source: template_fallback`), enabling secure, cost-free local evaluation.", C_GREEN),
    ]

    for i, (title, desc, color) in enumerate(narrative_steps):
        top_pos = Inches(1.4 + i * 1.8)
        add_card(s7, Inches(0.8), top_pos, Inches(11.7), Inches(1.6), bg_color=C_WHITE, border_color=C_BORDER)

        bar = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top_pos, Inches(0.2), Inches(1.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        tb = s7.shapes.add_textbox(Inches(1.2), top_pos + Inches(0.15), Inches(11.0), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color

        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_MUTED


    # ==========================================
    # SLIDE 8: MODEL PERFORMANCE & METRICS
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8, C_SLATE_BG)
    add_header(s8, "Model Performance, Validation & Test Coverage", "7. Empirical Validation")

    # 4 Metric KPI Cards
    metrics = [
        ("93.07%", "Model Accuracy", "XGBoost on 20% held-out test split", C_PRIMARY),
        ("0.9791", "ROC-AUC Score", "Superior discriminative capability", C_PURPLE),
        ("0.8526", "F1 Score", "Harmonic mean of precision & recall", C_GREEN),
        ("0.0486", "Brier Loss", "Well-calibrated probability estimates", C_NAVY_DARK),
    ]

    for i, (val, label, sub, col) in enumerate(metrics):
        left_pos = Inches(0.8 + i * 3.0)
        add_card(s8, left_pos, Inches(1.4), Inches(2.8), Inches(2.0), bg_color=C_WHITE, border_color=C_BORDER)

        tb = s8.shapes.add_textbox(left_pos + Inches(0.1), Inches(1.5), Inches(2.6), Inches(1.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = col
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = C_TEXT_DARK
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = sub
        p3.font.size = Pt(9)
        p3.font.color.rgb = C_TEXT_MUTED
        p3.alignment = PP_ALIGN.CENTER

    # Bottom: Test Coverage & Verification Details
    add_card(s8, Inches(0.8), Inches(3.7), Inches(11.7), Inches(3.2), bg_color=C_WHITE, border_color=C_BORDER)
    tb_t = s8.shapes.add_textbox(Inches(1.1), Inches(3.9), Inches(11.1), Inches(2.8))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    p = tf_t.paragraphs[0]
    p.text = "Comprehensive Automated Test Suite (`pytest backend/tests/ -v`)"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = C_PRIMARY

    test_items = [
        ("Thin-File Protection Guarantee:", "Verified with `test_thin_file_never_auto_declines` ensuring zero-score applicants are never rejected solely on missing credit history."),
        ("Behavioral Anomaly Escalation:", "Verified that extreme transaction-to-balance spikes escalate borderline approvals to human `REVIEW`."),
        ("Security & Auth Bounds:", "Role-scoped test assertions verifying 401 on missing keys and 403 on insufficient privilege scopes (Applicant vs Admin)."),
        ("End-to-End Decision Consistency:", "Automated regression tests covering all multi-modal fusion branches across synthetic customer archetypes.")
    ]
    for t, d in test_items:
        p = tf_t.add_paragraph()
        p.space_before = Pt(6)
        p.text = f"•  {t} {d}"
        p.font.size = Pt(11)
        p.font.color.rgb = C_TEXT_DARK


    # ==========================================
    # SLIDE 9: ENGINEERING DISCIPLINE & SECURITY
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9, C_SLATE_BG)
    add_header(s9, "Engineering Discipline, Cloud Readiness & Security", "8. Production Standards")

    eng_cards = [
        ("Security & Auth", "• Role-scoped API keys (Applicant vs Admin)\n• Strict Pydantic request schema validation\n• No hardcoded secrets (env configuration)\n• Non-root Docker execution", C_PRIMARY),
        ("Clean Architecture", "• Modular domain services (Credit, Behavior, Fusion, LLM)\n• API-First design with OpenAPI autodocs\n• Git version control with clean commit history\n• Structured error responses & logs", C_PURPLE),
        ("Containerization", "• Multi-stage Dockerfiles for backend and frontend\n• Single command launch: `docker compose up --build`\n• Isolated network bridging\n• Production-ready Nginx frontend bundle", C_GREEN),
        ("Cloud Blueprint", "• AWS ECS/Fargate container orchestration\n• AWS Bedrock / OpenAI LLM integration\n• PostgreSQL + pgvector for semantic search\n• CloudWatch telemetry & audit trails", C_NAVY_DARK),
    ]

    for i, (title, content, color) in enumerate(eng_cards):
        left_pos = Inches(0.8 + i * 3.0)
        add_card(s9, left_pos, Inches(1.4), Inches(2.8), Inches(5.5), bg_color=C_WHITE, border_color=C_BORDER)

        head = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, Inches(1.4), Inches(2.8), Inches(0.5))
        head.fill.solid()
        head.fill.fore_color.rgb = color
        head.line.fill.background()
        tf_h = head.text_frame
        p = tf_h.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = C_WHITE
        p.alignment = PP_ALIGN.CENTER

        tb = s9.shapes.add_textbox(left_pos + Inches(0.15), Inches(2.1), Inches(2.5), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(11)
        p.font.color.rgb = C_TEXT_DARK


    # ==========================================
    # SLIDE 10: ROADMAP & FUTURE VISION
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10, C_SLATE_BG)
    add_header(s10, "Future Roadmap: Scaling to Enterprise Production", "9. Vision & Roadmap")

    phases = [
        ("Phase 1: Semantic Underwriting Search (pgvector)", "Index historical credit applications and underwriter notes into PostgreSQL pgvector embeddings. Allows loan officers to query 'Find similar NTC gig-worker applications approved in Q3' for contextual precedence.", C_PRIMARY_BG, C_PRIMARY),
        ("Phase 2: Open Banking / Account Aggregator (AA)", "Connect real-time Open Banking APIs (e.g. Plaid / India Stack AA) to ingest continuous cashflow, GST filings, and merchant receipts directly into behavioral baselines.", C_PRIMARY_BG, C_PURPLE),
        ("Phase 3: Autonomous Underwriting Agents", "Deploy LangGraph / Bedrock multi-agent workflows that can autonomously request clarifying documents from applicants and cross-verify income before human final signoff.", C_GREEN_BG, C_GREEN),
    ]

    for i, (title, desc, bg, color) in enumerate(phases):
        top_pos = Inches(1.4 + i * 1.8)
        add_card(s10, Inches(0.8), top_pos, Inches(11.7), Inches(1.6), bg_color=C_WHITE, border_color=C_BORDER)

        bar = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top_pos, Inches(0.2), Inches(1.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        tb = s10.shapes.add_textbox(Inches(1.2), top_pos + Inches(0.15), Inches(11.0), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color

        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = C_TEXT_MUTED


    # ==========================================
    # SLIDE 11: CONCLUSION & CAMPUS PITCH SUMMARY
    # ==========================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_background(s11, C_NAVY_DARK)

    dec_card = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    dec_card.fill.solid()
    dec_card.fill.fore_color.rgb = C_NAVY_CARD
    dec_card.line.color.rgb = RGBColor(51, 65, 85)
    dec_card.line.width = Pt(1.5)

    tb = s11.shapes.add_textbox(Inches(1.3), Inches(1.2), Inches(10.7), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Why CredMe Stands Out"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = C_WHITE

    concl_points = [
        ("Solves the Exact Problem Statement:", "Purpose-built multi-modal engine expanding credit access to thin-file applicants while detecting live fraud."),
        ("High-Discipline Engineering:", "Clean modular FastAPI code, Pydantic validation, role-based auth, Dockerized environment, and automated test suite."),
        ("Responsible & Explainable AI:", "Mathematical SHAP feature attributions, disparate impact audits, and guardrailed generative explanations."),
        ("Working End-to-End Prototype:", "Interactive React 19 UI, live WebSocket streaming, and verified reproducible setup in < 5 minutes.")
    ]

    for title, desc in concl_points:
        p = tf.add_paragraph()
        p.space_before = Pt(12)
        p.text = f"✓  {title} "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(165, 180, 252)
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(203, 213, 225)

    p_end = tf.add_paragraph()
    p_end.space_before = Pt(18)
    p_end.text = "Thank you! Ready for live demonstration & pitch during campus evaluation."
    p_end.font.size = Pt(14)
    p_end.font.bold = True
    p_end.font.color.rgb = C_GREEN

    prs.save(output_filename)
    print(f"Presentation saved successfully to {output_filename}")

if __name__ == "__main__":
    create_presentation("CredMe-Pitch-Deck.pptx")
    create_presentation("CredMe-status-deck.pptx")
