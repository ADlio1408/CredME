from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Next-Gen Credit Intelligence: CredMe Status"
subtitle.text = "Summary, Architecture, Work Completed, and Next Steps"

# Slide 2: Objective + Challenge
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Objective & Challenge"
body = slide.shapes.placeholders[1].text_frame
body.text = "Objective: Expand credit access to New-to-Credit (NTC) and thin-file customers using alternative data and real-time behavioral signals."

p = body.add_paragraph()
p.text = "Challenge: Traditional models rely on static historical data and miss many creditworthy individuals. Shift towards proactive, contextual decisioning."

# Slide 3: Completed Work (Commits)
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Completed Work (Recent Commits)"
body = slide.shapes.placeholders[1].text_frame
body.text = "- Git setup, README, Docker (E2E verified)"
body.add_paragraph().text = "- PaymentHistory bug fix"
body.add_paragraph().text = "- Fairness audit and report"
body.add_paragraph().text = "- Role-scoped API keys (Applicant / Admin)"
body.add_paragraph().text = "- Real-time transaction streaming + WebSocket"
body.add_paragraph().text = "- LLM narrative scaffold (template fallback)"

# Slide 4: Current Work - RentPaymentConsistency
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "In Progress: Alternative Data Signal"
body = slide.shapes.placeholders[1].text_frame
body.text = "Adding illustrative signal: RentPaymentConsistency"
body.add_paragraph().text = "- Integrated into Decision API for thin-file applicants"
body.add_paragraph().text = "- Threshold: flag if < 0.70 as financial concern"
body.add_paragraph().text = "- Unit tests added to assert behavior"

# Slide 5: Architecture Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "High-level Architecture"
body = slide.shapes.placeholders[1].text_frame
body.text = "Frontend: ReactJS (Vite)"
body.add_paragraph().text = "Backend: FastAPI (Python) — decision & behavior engines"
body.add_paragraph().text = "DB: PostgreSQL (planned pgvector for semantic search)"
body.add_paragraph().text = "AI Layer: LLM via configurable API (template fallback if no key)"
body.add_paragraph().text = "Streaming: WebSocket endpoint for live transaction ingestion"

# Slide 6: Next Steps
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Next Steps"
body = slide.shapes.placeholders[1].text_frame
body.text = "- Complete RentPaymentConsistency integration and run full test suite"
body.add_paragraph().text = "- Implement pgvector + semantic search endpoint (/semantic/similar-cases)"
body.add_paragraph().text = "- Add monitoring, secure secret management, and CI tests"
body.add_paragraph().text = "- Prepare demo and submission PPT/PDF per hackathon guidelines"

# Save
prs.save('CredMe-status-deck.pptx')
print('Created CredMe-status-deck.pptx')
